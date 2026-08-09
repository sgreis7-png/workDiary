"""Build the onboarding deck as a real .pptx, from the same screenshots as the web version.

python-pptx has no RTL switch, so every paragraph gets `rtl="1"` set on its XML directly
and is right-aligned — without it Hebrew punctuation and mixed Hebrew/Latin lines
(".mpp", "MS Project") come out in the wrong order.
"""
import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from build_deck import IMAGES, encode  # noqa: F401  (IMAGES is the shared source list)
from PIL import Image

OUT = pathlib.Path(__file__).parent / "אגרוטופ — היכרות עם המערכת.pptx"
TMP = pathlib.Path(__file__).parent / ".pptx-cache"
TMP.mkdir(exist_ok=True)

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.62)

INK = RGBColor(0x0B, 0x12, 0x10)
INK2 = RGBColor(0x26, 0x31, 0x2C)
FAINT = RGBColor(0x68, 0x76, 0x6F)
GREEN = RGBColor(0x00, 0x85, 0x40)
FOREST = RGBColor(0x0D, 0x1F, 0x15)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xEE, 0xF2, 0xF0)
LIGHT = RGBColor(0xE9, 0xF3, 0xEE)
DIM = RGBColor(0xCF, 0xE0, 0xD6)
CLAY = RGBColor(0xC1, 0x4A, 0x15)

FONT = "Segoe UI"


def prepared(key: str) -> pathlib.Path:
    """Resize the source screenshot once and cache it next to this script."""
    src, max_w, trim = IMAGES[key]
    dst = TMP / f"{key}.png"
    if dst.exists():
        return dst
    im = Image.open(src)
    if trim:
        from PIL import ImageChops
        rgb = im.convert("RGB")
        box = ImageChops.difference(rgb, Image.new("RGB", rgb.size, (255, 255, 255))).convert("L").getbbox()
        if box:
            im = im.crop(box)
    if im.width > max_w:
        im = im.resize((max_w, round(im.height * max_w / im.width)), Image.LANCZOS)
    im.convert("RGB").save(dst, format="PNG")
    return dst


def rtl(paragraph, align=PP_ALIGN.RIGHT) -> None:
    """Mark the paragraph right-to-left on the underlying <a:pPr>.

    python-pptx exposes no RTL property, and without the attribute PowerPoint lays Hebrew
    out left-to-right: trailing punctuation jumps to the wrong end and mixed lines such as
    "MS Project", ".mpp" or "שפות, RTL מלא" come out reversed. Centred text needs the
    attribute just as much as right-aligned text, hence the overridable alignment.
    """
    paragraph._p.get_or_add_pPr().set("rtl", "1")
    paragraph.alignment = align


def bg(slide, colour) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def text(slide, body, *, left, top, width, height, size, colour, bold=False,
         line=1.25, space_after=6, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = body if isinstance(body, list) else [body]
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        bullet, content = (True, item[1]) if isinstance(item, tuple) else (False, item)
        run = p.add_run()
        run.text = ("•  " if bullet else "") + content
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = colour
        run.font.name = FONT
        p.line_spacing = line
        p.space_after = Pt(space_after)
        rtl(p)
    return box


def picture(slide, key, *, top, bottom_margin=Inches(0.45)):
    """Fit the screenshot into the space left below `top`, centred."""
    path = prepared(key)
    iw, ih = Image.open(path).size
    avail_w = W - 2 * MARGIN
    avail_h = H - top - bottom_margin
    scale = min(avail_w / iw, avail_h / ih)
    w, h = Emu(int(iw * scale)), Emu(int(ih * scale))
    return slide.shapes.add_picture(str(path), Emu(int((W - w) / 2)), top, width=w, height=h)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def header(slide, eyebrow, title, *, dark=False, body=None):
    """Eyebrow + title (+ optional lead), returning the y where content may start."""
    ec, tc, bc = (RGBColor(0x8F, 0xB7, 0xA0), PAPER, DIM) if dark else (FAINT, INK, INK2)
    text(slide, eyebrow, left=MARGIN, top=Inches(0.42), width=W - 2 * MARGIN, height=Inches(0.3),
         size=11, colour=ec, bold=True)
    text(slide, title, left=MARGIN, top=Inches(0.72), width=W - 2 * MARGIN, height=Inches(0.95),
         size=30, colour=tc, bold=True, line=1.1)
    y = Inches(1.78)
    if body:
        n = len(body) if isinstance(body, list) else 1
        h = Inches(0.32) * n + Inches(0.18)
        text(slide, body, left=MARGIN, top=y, width=W - 2 * MARGIN, height=h,
             size=14, colour=bc, line=1.35)
        y = y + h + Inches(0.12)
    return y


def cards(slide, items, *, top, cols=3, dark=False):
    """A row (or rows) of titled cards, drawn as rounded rectangles."""
    from pptx.enum.shapes import MSO_SHAPE
    gap = Inches(0.22)
    total = W - 2 * MARGIN
    cw = Emu(int((total - gap * (cols - 1)) / cols))
    ch = Inches(1.62)
    for i, (head, para) in enumerate(items):
        row, col = divmod(i, cols)
        x = W - MARGIN - cw - (cw + gap) * col          # RTL: first card on the right
        y = top + (ch + gap) * row
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x18, 0x36, 0x27) if dark else PAPER
        shape.line.color.rgb = RGBColor(0x2A, 0x4A, 0x38) if dark else RGBColor(0xD7, 0xDF, 0xDB)
        shape.line.width = Pt(0.75)
        shape.shadow.inherit = False
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.16)
        tf.margin_top = tf.margin_bottom = Inches(0.12)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = head
        r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = PAPER if dark else INK
        r.font.name = FONT
        rtl(p)
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = para
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = DIM if dark else INK2
        r2.font.name = FONT
        p2.line_spacing = 1.25
        rtl(p2)
    return top + (ch + gap) * ((len(items) + cols - 1) // cols)


def stats(slide, items, *, top, dark=False):
    from pptx.enum.shapes import MSO_SHAPE
    gap = Inches(0.2)
    cols = len(items)
    total = W - 2 * MARGIN
    cw = Emu(int((total - gap * (cols - 1)) / cols))
    ch = Inches(1.05)
    for i, (value, label) in enumerate(items):
        x = W - MARGIN - cw - (cw + gap) * i
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, cw, ch)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x18, 0x36, 0x27) if dark else PAPER
        shape.line.color.rgb = RGBColor(0x2A, 0x4A, 0x38) if dark else RGBColor(0xD7, 0xDF, 0xDB)
        shape.line.width = Pt(0.75)
        shape.shadow.inherit = False
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = value
        r.font.size = Pt(24); r.font.bold = True
        r.font.color.rgb = PAPER if dark else INK
        r.font.name = FONT
        rtl(p, PP_ALIGN.CENTER)
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = label
        r2.font.size = Pt(10)
        r2.font.color.rgb = DIM if dark else FAINT
        r2.font.name = FONT
        rtl(p2, PP_ALIGN.CENTER)
    return top + ch


def table_slide(slide, rows, *, top):
    cols = 3
    height = Inches(0.42) * (len(rows) + 1)
    shape = slide.shapes.add_table(len(rows) + 1, cols, MARGIN, top, W - 2 * MARGIN, height)
    tbl = shape.table
    tbl.columns[0].width = Emu(int((W - 2 * MARGIN) * 0.40))   # RTL: first column renders right
    tbl.columns[1].width = Emu(int((W - 2 * MARGIN) * 0.38))
    tbl.columns[2].width = Emu(int((W - 2 * MARGIN) * 0.22))
    head = ['איך זה עובד עכשיו', 'איך זה היה', 'המשימה']
    for c, label in enumerate(head):
        cell = tbl.cell(0, c)
        cell.text = ''
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = label
        r.font.size = Pt(11); r.font.bold = True
        r.font.color.rgb = FAINT; r.font.name = FONT
        rtl(p)
    for i, (task, old, new) in enumerate(rows, start=1):
        for c, value in enumerate([new, old, task]):
            cell = tbl.cell(i, c)
            cell.text = ''
            p = cell.text_frame.paragraphs[0]
            r = p.add_run(); r.text = value
            r.font.size = Pt(10.5)
            r.font.bold = (c == 2)
            r.font.color.rgb = INK if c == 0 else INK2
            r.font.name = FONT
            rtl(p)


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ---- 01 cover ----
    # Pure white, not cream: the logo file carries its own white background, which would
    # otherwise read as a pale rectangle floating on the slide.
    s = blank(prs); bg(s, PAPER)
    logo = prepared('logo')
    lw = Inches(3.1)
    liw, lih = Image.open(logo).size
    s.shapes.add_picture(str(logo), Emu(int(W - MARGIN - lw)), Inches(0.6),
                         width=lw, height=Emu(int(lw * lih / liw)))
    text(s, 'היכרות עם המערכת · לעובדי אגרוטופ', left=MARGIN, top=Inches(2.5),
         width=W - 2 * MARGIN, height=Inches(0.34), size=12, colour=FAINT, bold=True)
    text(s, 'יומן עבודה ומרכז בקרה לפרויקט', left=MARGIN, top=Inches(2.85),
         width=W - 2 * MARGIN, height=Inches(1.5), size=44, colour=INK, bold=True, line=1.08)
    text(s, 'כל מה שקורה בפרויקט — דיווח יומי מהשטח, לוח זמנים, בקרת איכות וליקויים, '
            'התקדמות ועובדים — במקום אחד, מכל מחשב וכל טלפון.',
         left=MARGIN, top=Inches(4.35), width=Inches(8.4), height=Inches(1.1),
         size=15, colour=INK2, line=1.4)
    text(s, 'work-diary-phi.vercel.app   ·   עברית · English   ·   דפדפן · אנדרואיד · אופליין',
         left=MARGIN, top=Inches(6.45), width=W - 2 * MARGIN, height=Inches(0.4),
         size=11, colour=FAINT)

    # ---- 02 the problem ----
    s = blank(prs); bg(s, CREAM)
    y = header(s, 'למה בכלל', 'איך זה עובד היום — ולמה זה נשבר',
               body='כל הכלים האלה טובים בדבר אחד. הבעיה היא שאף אחד מהם לא מחזיק את התמונה '
                    'המלאה, ואף אחד לא מדבר עם השני.')
    cards(s, [
        ('וואטסאפ', 'הדיווח מגיע מהר ונעלם מהר. אין חיפוש אמיתי ואין תמונת מצב.'),
        ('אקסל', 'גרסאות. מי עדכן מה, איזה קובץ נכון. עובד לאדם אחד, נשבר בצוות.'),
        ('MS Project', 'לוח זמנים מעולה — שיושב על מחשב אחד. מי שבשטח לא רואה אותו.'),
        ('דוחות נייר', 'צ׳קליסט, חתימות וליקויים בקלסר. אי אפשר לספור או להתריע.'),
        ('תמונות בגלריה', 'אלפי תמונות בלי הקשר — לאיזה לול, לאיזה תאריך, לאיזה ליקוי.'),
        ('הידע בראש', 'מי שמכיר את הפרויקט יוצא לחופשה, והפרויקט נעצר איתו.'),
    ], top=y, cols=3)

    # ---- 03 the answer ----
    s = blank(prs); bg(s, FOREST)
    y = header(s, 'הרעיון', 'מקור אמת אחד לכל פרויקט', dark=True,
               body='מדווחים פעם אחת — מהשטח, מהטלפון, גם בלי קליטה. משם הנתון עובד לבד: נכנס '
                    'לדוח, לגרף ההתקדמות, לסיכום השבועי, למרכז הבקרה ולהתראות.')
    y = stats(s, [('1', 'הזנה אחת'), ('6', 'שכבות בפרויקט'),
                  ('∞', 'היסטוריה שנשמרת'), ('0', 'קבצים לשלוח במייל')], top=y, dark=True)
    cards(s, [
        ('מה שהעובד עושה', 'ממלא רשומה יומית עם תמונות · מסמן התקדמות לכל לול · '
                           'מעדכן צ׳קליסט וליקויים'),
        ('מה שהמערכת עושה מזה', 'גרפי התקדמות לאורך זמן · סיכום שבועי עם תובנות · '
                                'התראות על חריגה מיעד'),
    ], top=y + Inches(0.3), cols=2, dark=True)

    # ---- screenshot slides ----
    shots = [
        ('כניסה', 'נכנסים מהדפדפן, בלי להתקין כלום',
         'הכניסה בדוא״ל של אגרוטופ, ורק כתובות שאושרו מראש יכולות להירשם. בטלפון מוסיפים '
         'למסך הבית וזה מתנהג כאפליקציה.', 'login'),
        ('התמצאות', 'התפריט — לפי נושאים, לא רשימה ארוכה',
         'שש קבוצות. הקבוצה שאתה נמצא בה נפתחת לבד והשאר מקופלות, כך שרואים חמישה־שישה '
         'דברים ולא ארבעה־עשר.', 'menu'),
        ('הליבה', 'הרשומה היומית — הבסיס לכל השאר',
         'מי היה בשטח, מה נעשה, כמה עובדים ובאילו שעות, מזג אוויר, תמונות, בלת״מ ובטיחות. '
         'אותה רשומה היא גם הדוח ללקוח וגם הנתון שמזין את הגרפים.', 'logbook'),
        ('התקדמות', 'גרף התקדמות לכל לול, לאורך זמן',
         'העובד מזיז מחוון אחוזים בדיווח היומי, ומזה נבנית עקומה. רואים מיד מי מתקדם, '
         'מי נעצר ומתי בדיוק זה קרה.', 'progress'),
        ('לוח זמנים', 'קובץ MS Project נכנס — ויוצא לוח חי',
         'מעלים את קובץ ה־mpp. כמו שהוא. הוא מומר מאחורי הקלעים ונפרס ללוח זמנים שכל מי '
         'שמורשה רואה — בלי להתקין MS Project ובלי לשלוח קבצים.', 'gantt'),
        ('לוח זמנים', 'גורר משימה — כל מה שתלוי בה נדחה איתה',
         'המערכת דוחפת אוטומטית את המשימות התלויות לפי סוג הקשר, מעדכנת את שורות הסיכום, '
         'ומראה בכמה ימים המשימה זזה מהתכנון המקורי.', 'gantt_hover'),
        ('בקרת איכות', 'תפיסת סיום שלב — הקלסר הופך למערכת',
         'שבעה שערים לכל לול. פריט שסומן «לא בוצע» הופך אוטומטית לליקוי עם חומרה, אחראי '
         'ותאריך יעד. אי אפשר לחתום על שער עם ליקוי קריטי פתוח.', 'qc_gate'),
        ('בקרת איכות', 'פתיחת לול — מי אחראי על מה, מוסכם מראש',
         'גז, גנרטור, קו מים, חשמל, ציוד גידול — לכל תחום נקבע אם הוא על אגרוטופ, על הלקוח '
         'או על גורם חוץ. כך אין ויכוח באמצע הדרך.', 'qc_open'),
        ('בקרת איכות', 'יומן ליקויים — עם שם, יעד וסטטוס',
         'לכל ליקוי מספר, חומרה, אחראי ותאריך יעד. איחור לא נשאר שקט: התראה לאחראי, '
         'והסלמה למנהלים כשהאיחור מתארך.', 'qc_log'),
        ('בקרת איכות', 'ריכוז סטטוס לכל שער — באחוזים',
         'אותו חישוב שהיה בגיליון האקסל, רק שהוא מתעדכן לבד: בוצע, לא בוצע, לא רלוונטי, '
         'טרם — ואחוז לכל שער.', 'qc_summary'),
        ('פלט', 'דוח מוכן לשליחה — בלחיצה',
         'אותם נתונים בפריסת דוח: פרטי הלול, ריכוז הסטטוס, הליקויים הפתוחים והחתימות. '
         'אין הקלדה מחדש, ולכן אין פער בין הדוח למה שבשטח.', 'qc_report'),
        ('מעקב', 'סיכום שבועי — נשלח לבד בבוקר ראשון',
         'כרטיס לכל פרויקט: מה נעשה השבוע, התקדמות לכל לול עם הפרש מהשבוע שעבר, ליקויים '
         'שנפתחו ונסגרו, בטיחות — ופרויקטים שלא דיווחו בכלל.', 'digest'),
        ('למנהלים', 'מרכז בקרה — כל הפרויקט במסך אחד',
         'למעלה המספרים שקובעים. מתחת שש שכבות — סקירה, לוח זמנים, לולים, ליקויים, עובדים '
         'ויומן — לפי טאבים או הכל בעמוד אחד. לאדמין ולמנהלים בלבד.', 'control'),
    ]
    for eyebrow, title, body, key in shots:
        s = blank(prs); bg(s, CREAM)
        y = header(s, eyebrow, title, body=body)
        picture(s, key, top=y)

    # ---- phone slide ----
    s = blank(prs); bg(s, CREAM)
    header(s, 'בשטח', 'בטלפון — לצפייה, בכוונה')
    path = prepared('gantt_phone')
    iw, ih = Image.open(path).size
    ph = Inches(4.6)
    pw = Emu(int(ph * iw / ih))
    s.shapes.add_picture(str(path), Emu(int(MARGIN + Inches(0.5))), Inches(2.1), width=pw, height=ph)
    text(s, ['הלוח נפתח בטלפון מוקטן לכל הפרויקט, עם שמות המשימות, אבני הדרך והקו של היום. '
             'לחיצה על שורה פותחת את פרטיה מעל הלוח.',
             'עריכה חסומה בטלפון במכוון: בזום כזה רוחב הבר הוא כמה פיקסלים, וגרירה בטעות '
             'הייתה מזיזה את המשימה הלא נכונה. מכיוון שהשינוי נשמר לכולם ודוחף את כל '
             'התלויות, טעות אחת מספיקה.',
             'לשינוי תאריכים פותחים במחשב. אותו מידע, פחות דרכים לטעות.'],
         left=Emu(int(MARGIN + pw + Inches(1.1))), top=Inches(2.3),
         width=Emu(int(W - 2 * MARGIN - pw - Inches(1.1))), height=Inches(4),
         size=14, colour=INK2, line=1.4, space_after=12)

    # ---- comparison ----
    s = blank(prs); bg(s, CREAM)
    y = header(s, 'ההשוואה', 'מה משתנה בפועל')
    table_slide(s, [
        ('דיווח יומי', 'הודעה בוואטסאפ, נעלמת בגלילה', 'רשומה עם תמונות, נשמרת וניתנת לחיפוש'),
        ('מצב לוח הזמנים', 'קובץ על מחשב אחד', 'לוח חי, זמין לכל מי שמורשה, גם בטלפון'),
        ('שינוי תאריך', 'מעדכנים ומקווים שכולם ראו', 'התלויות נדחות לבד, השינוי מתועד'),
        ('התקדמות', 'מרכיבים אקסל, שואלים אנשים', 'גרף לכל לול, מתעדכן מהדיווח'),
        ('סיום שלב', 'צ׳קליסט בקלסר', 'שערים עם אכיפה, חתימות ותמונות'),
        ('ליקוי שנשכח', 'מתגלה במסירה', 'התראה לאחראי, הסלמה למנהל'),
        ('דוח ללקוח', 'מקלידים מחדש בוורד', 'מופק מהנתונים, בלחיצה'),
        ('אין קליטה בשטח', 'רושמים על דף ומקווים לזכור', 'נשמר במכשיר ונשלח כשחוזרת רשת'),
    ], top=y)

    # ---- serious ----
    s = blank(prs); bg(s, FOREST)
    y = header(s, 'מתחת למכסה', 'למה זו מערכת ולא תחביב', dark=True,
               body='הדברים שלא רואים במסך הם אלה שקובעים אם אפשר לסמוך על המערכת עם נתוני '
                    'פרויקט אמיתיים.')
    y = cards(s, [
        ('אבטחה', 'הרשאה נאכפת במסד הנתונים ולא במסך · הרשאות לפי תחום לכל משתמש · '
                  'הרשמה רק לכתובות שאושרו · קבצים בקישורים חתומים לשעה'),
        ('אמינות', 'יומן שינויים לפעולות רגישות · תור אופליין שנשלח לבד · '
                   'בדיקות אוטומטיות על כל שינוי בקוד · גרסאות מסד נתונים מתועדות'),
    ], top=y, cols=2, dark=True)
    stats(s, [('180', 'בדיקות אוטומטיות'), ('44', 'גרסאות מסד נתונים'),
              ('10', 'תחומי הרשאה'), ('2', 'שפות, RTL מלא')], top=y + Inches(0.3), dark=True)

    # ---- getting started ----
    s = blank(prs); bg(s, CREAM)
    y = header(s, 'להתחיל', 'שלושה דברים, וסיימת')
    y = cards(s, [
        ('① נכנסים ומוסיפים למסך הבית', 'משם זה נראה ומתנהג כאפליקציה, כולל עבודה בלי קליטה.'),
        ('② מדווחים בסוף כל יום', 'רשומה אחת ליום: מה נעשה, כמה אנשים, תמונות, ואחוז לכל לול.'),
        ('③ מאשרים התראות', 'כדי לקבל הודעה על ליקוי שהוקצה לך ועל תאריך יעד שמתקרב.'),
    ], top=y, cols=3)
    text(s, ['שאלות נפוצות',
             'אני לא רואה תפריט שמישהו אחר רואה — ההרשאות אישיות לפי תחום. פנה לאדמין.',
             'דיווחתי בלי קליטה, זה נעלם? לא. נשמר במכשיר ונשלח כשחוזרת רשת.',
             'אני רואה גרסה ישנה של המסך — רענון מלא: Ctrl+Shift+R, ובטלפון גלילה מלמעלה למטה.'],
         left=MARGIN, top=y + Inches(0.35), width=W - 2 * MARGIN, height=Inches(1.7),
         size=12.5, colour=INK2, line=1.4)

    # ---- closing ----
    s = blank(prs); bg(s, FOREST)
    text(s, 'מדווחים פעם אחת.', left=MARGIN, top=Inches(2.5), width=W - 2 * MARGIN,
         height=Inches(0.95), size=40, colour=PAPER, bold=True)
    text(s, 'המערכת עושה את השאר.', left=MARGIN, top=Inches(3.4), width=W - 2 * MARGIN,
         height=Inches(0.95), size=40, colour=RGBColor(0x6F, 0xD0, 0x6A), bold=True)
    text(s, 'עדיף לשאול מאשר לנחש — לכל שאלה, בעיה או רעיון לשיפור יש כפתור «משוב» '
            'בתחתית התפריט.',
         left=MARGIN, top=Inches(4.6), width=Inches(8), height=Inches(0.9),
         size=15, colour=DIM, line=1.4)
    text(s, 'work-diary-phi.vercel.app', left=MARGIN, top=Inches(6.4),
         width=W - 2 * MARGIN, height=Inches(0.4), size=11, colour=RGBColor(0x8F, 0xB7, 0xA0))

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    # the console here is cp1252; a Hebrew path in a print would raise after a good save
    print(f"saved {len(prs.slides._sldIdLst)} slides, "
          f"{OUT.stat().st_size / 1024 / 1024:.1f} MB, beside this script")


if __name__ == "__main__":
    build()
